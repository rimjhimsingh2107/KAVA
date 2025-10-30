#!/usr/bin/env python3
"""
Generate KAVA Project Presentation
Professional design with clean visuals
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (professional blue/orange)
PRIMARY_COLOR = RGBColor(26, 54, 93)      # Dark blue
ACCENT_COLOR = RGBColor(255, 127, 14)     # Orange
LIGHT_BG = RGBColor(240, 244, 248)        # Light blue
SUCCESS_COLOR = RGBColor(34, 139, 34)     # Green
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Title slide with gradient effect"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.alignment = PP_ALIGN.CENTER
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.alignment = PP_ALIGN.CENTER
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = ACCENT_COLOR
    
    return slide

def add_content_slide(prs, title, content_items):
    """Content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.space_before = Pt(8)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_items, right_items):
    """Two-column layout slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title bar
    title_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = PRIMARY_COLOR
    title_shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5), Inches(4), Inches(5.5)
    )
    left_frame = left_box.text_frame
    for i, item in enumerate(left_items):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.2), Inches(1.5), Inches(4), Inches(5.5)
    )
    right_frame = right_box.text_frame
    for i, item in enumerate(right_items):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(6)
    
    return slide

# SLIDE 1: Title
add_title_slide(prs, "KAVA", "AI-Powered Insurance Claims Validation")

# SLIDE 2: The Problem
add_content_slide(prs, "The Problem", [
    "⏱️ Insurance claim processing takes 2-6 weeks on average",
    "📄 Manual document review is time-consuming and error-prone",
    "💰 Fraud costs the insurance industry $80+ billion annually",
    "🔍 Claims adjusters spend hours validating receipts and evidence",
    "❌ 30% of claims require multiple resubmissions due to missing docs"
])

# SLIDE 3: The Solution
add_content_slide(prs, "The Solution: KAVA", [
    "🤖 AI-powered validation using Claude Sonnet 4.5",
    "⚡ Processes complete claims in under 2 minutes",
    "📊 47-rule constitution for comprehensive analysis",
    "🔄 Progressive improvement loop with auto-enhancement",
    "🔐 Blockchain-ready cryptographic attestation",
    "📦 Generates submission-ready insurance packages"
])

# SLIDE 4: System Architecture
add_two_column_slide(prs, "System Architecture", [
    "🎨 FRONTEND",
    "• Next.js 15 + TypeScript",
    "• 4-step wizard interface",
    "• Real-time validation feedback",
    "• Port 3001",
    "",
    "🔧 BACKEND",
    "• FastAPI (Python)",
    "• Claude AI integration",
    "• SQLite database",
    "• Port 8000"
], [
    "🤖 AI SERVICES",
    "• Claude Vision API (OCR)",
    "• AI Judge (validation)",
    "• Document processor",
    "• Receipt auto-fetch (Knot API)",
    "",
    "📊 OUTPUTS",
    "• PDF generation (ReportLab)",
    "• ZIP packaging",
    "• ECDSA signatures",
    "• Blockchain proofs"
])

# SLIDE 5: The 4-Step Flow
add_content_slide(prs, "User Journey: 4 Simple Steps", [
    "1️⃣ UPLOAD DOCUMENTS",
    "   Photos, receipts, policy, fire reports → Claude OCR extracts data",
    "",
    "2️⃣ CREATE CLAIM PACKET",
    "   Fill in basic info → Generate organized claim packet PDF",
    "",
    "3️⃣ AI JUDGE VALIDATION",
    "   Progressive 4-iteration loop → Auto-improvement → 47-rule analysis",
    "",
    "4️⃣ GENERATE OUTPUTS",
    "   Complete ZIP package + Validation report + Cryptographic proof"
])

# SLIDE 6: AI Judge - The Brain
add_content_slide(prs, "AI Judge: Progressive Validation Loop", [
    "📊 47-Rule Constitution across 7 categories:",
    "   • Completeness (12 rules, 45% weight)",
    "   • Damage Assessment (10 rules, 28% weight)",
    "   • Documentation Quality (6 rules, 27% weight)",
    "   • + Temporal, Geographic, Policy, Financial",
    "",
    "🔄 4 Progressive Iterations:",
    "   Iteration 1: Baseline screening (53% typical)",
    "   Iteration 2: + Knot receipts auto-fetch (74% typical)",
    "   Iteration 3: + Deep OCR reprocessing (81% typical)",
    "   Iteration 4: Final expert review (85% typical)",
    "",
    "🎯 Target: 80%+ score = automatic approval"
])

# SLIDE 7: Live Demo Results
add_two_column_slide(prs, "Demo: Real Wildfire Claim", [
    "📥 DOCUMENTS UPLOADED:",
    "• Kitchen fire damage photo",
    "• Home Depot receipt ($1,073)",
    "• Lowe's receipt ($289)",
    "• Insurance policy",
    "• Fire dept report",
    "• Contractor estimate ($115k)",
    "",
    "⏱️ PROCESSING TIME:",
    "• Document OCR: 8 seconds",
    "• Validation loop: 47 seconds",
    "• Output generation: 12 seconds",
    "• Total: 67 seconds"
], [
    "📊 VALIDATION RESULTS:",
    "• Iteration 1: 25/47 rules = 53%",
    "• Iteration 2: 35/47 rules = 74%",
    "  (+ 6 Knot receipts fetched)",
    "• Iteration 3: 38/47 rules = 81%",
    "  ✅ Target reached!",
    "",
    "🏆 FINAL OUTPUTS:",
    "• 9 professional PDFs in ZIP",
    "• $116,713 itemized inventory",
    "• 🥈 SILVER_TRUST badge",
    "• ECDSA cryptographic proof"
])

# SLIDE 8: Technical Innovation
add_content_slide(prs, "Technical Innovation", [
    "🧠 Real AI - Not Mock Data",
    "   Claude Sonnet 4.5 for OCR and validation (not hardcoded responses)",
    "",
    "🔄 Progressive Improvement",
    "   System actively enhances claims by auto-fetching receipts",
    "",
    "📏 47-Rule Constitution",
    "   Comprehensive validation covering all aspects of wildfire claims",
    "",
    "🔐 Cryptographic Attestation",
    "   ECDSA signatures (SHA-256) for tamper-proof validation proof",
    "",
    "💾 Production-Ready",
    "   Database persistence, error handling, full audit trail"
])

# SLIDE 9: Key Algorithms
add_two_column_slide(prs, "Key Algorithms & Logic", [
    "📊 SCORE CALCULATION:",
    "Rules Passed / Total Rules",
    "",
    "Example:",
    "38 rules passed out of 47",
    "= 38/47 = 80.85%",
    "",
    "Higher score = more rules",
    "satisfied by the claim",
    "",
    "Each rule weighted by",
    "importance (5%-20%)"
], [
    "💰 VALUE EXTRACTION:",
    "Receipt OCR gets amounts:",
    "[899, 24.99, 18.97, 1073.86]",
    "         ↑",
    "   Takes last value (total)",
    "",
    "Itemized Inventory:",
    "Sum of all receipt totals",
    "$1,073 + $289 + $115,350",
    "= $116,713.24",
    "",
    "All calculated automatically!"
])

# SLIDE 10: Generated Outputs
add_content_slide(prs, "Complete Claim Package (9 PDFs)", [
    "📦 COMPREHENSIVE ZIP PACKAGE:",
    "",
    "1. 📄 Cover Letter - Professional submission letter",
    "2. 📋 Proof of Loss Statement - Official sworn claim",
    "3. 📸 Property Damage Photos - Embedded images with descriptions",
    "4. 📊 Itemized Loss Inventory - $116,713.24 documented",
    "5. 🧾 Purchase Receipts - All receipts compiled",
    "6. 🚒 Fire Department Report - Official incident documentation",
    "7. 🔨 Contractor Estimates - Professional assessments",
    "8. 🤖 AI Validation Report - 38/47 rules passed, detailed rationale",
    "9. 🔐 Cryptographic Proof Card - ECDSA signature for verification"
])

# SLIDE 11: Tech Stack
add_two_column_slide(prs, "Technology Stack", [
    "🎨 FRONTEND:",
    "• Next.js 15",
    "• TypeScript",
    "• Tailwind CSS",
    "• Framer Motion",
    "• React Dropzone",
    "",
    "🔧 BACKEND:",
    "• FastAPI (Python)",
    "• Claude AI API",
    "• SQLAlchemy",
    "• SQLite"
], [
    "🤖 AI & SERVICES:",
    "• Claude Sonnet 4.5",
    "• Computer Vision OCR",
    "• Knot API (receipts)",
    "• PyPDF2 (text extraction)",
    "",
    "📄 PDF & CRYPTO:",
    "• ReportLab (PDF gen)",
    "• Cryptography lib",
    "• ECDSA signatures",
    "• SHA-256 hashing"
])

# SLIDE 12: Impact & Value
add_content_slide(prs, "Business Impact", [
    "⚡ 95% faster processing (weeks → minutes)",
    "",
    "💰 Reduces fraud risk with automated pattern detection",
    "",
    "🎯 80%+ accuracy with AI-powered validation",
    "",
    "📉 Decreases resubmission rate by auto-identifying missing docs",
    "",
    "✅ Improves customer experience with instant feedback",
    "",
    "🔐 Provides blockchain-verifiable proof for compliance"
])

# SLIDE 13: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.shapes.add_shape(
    1, Inches(0), Inches(0), Inches(10), Inches(7.5)
)
background.fill.solid()
background.fill.fore_color.rgb = PRIMARY_COLOR
background.line.fill.background()

title_box = slide.shapes.add_textbox(
    Inches(1), Inches(2.5), Inches(8), Inches(1)
)
title_frame = title_box.text_frame
title_frame.text = "Thank You!"
title_p = title_frame.paragraphs[0]
title_p.alignment = PP_ALIGN.CENTER
title_p.font.size = Pt(60)
title_p.font.bold = True
title_p.font.color.rgb = WHITE

subtitle_box = slide.shapes.add_textbox(
    Inches(1), Inches(3.8), Inches(8), Inches(2)
)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Questions?"
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.alignment = PP_ALIGN.CENTER
subtitle_p.font.size = Pt(36)
subtitle_p.font.color.rgb = ACCENT_COLOR

contact_box = slide.shapes.add_textbox(
    Inches(2), Inches(5.5), Inches(6), Inches(1)
)
contact_frame = contact_box.text_frame
contact_frame.text = "GitHub: github.com/rimjhimsingh2107/KAVA"
contact_p = contact_frame.paragraphs[0]
contact_p.alignment = PP_ALIGN.CENTER
contact_p.font.size = Pt(18)
contact_p.font.color.rgb = WHITE

# Save presentation
output_path = "/Users/rimjhim/Desktop/SELFPROJECTS/Kava- testing/KAVA_Presentation.pptx"
prs.save(output_path)

print("✅ Presentation created successfully!")
print(f"📁 Location: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("\n🎯 Slides created:")
print("  1. Title Slide")
print("  2. The Problem")
print("  3. The Solution")
print("  4. System Architecture")
print("  5. User Journey")
print("  6. AI Judge Engine")
print("  7. Live Demo Results")
print("  8. Technical Innovation")
print("  9. Key Algorithms")
print(" 10. Generated Outputs")
print(" 11. Tech Stack")
print(" 12. Business Impact")
print(" 13. Thank You")
